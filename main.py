#!/usr/bin/env python3
from __future__ import annotations

import json
import re
import unicodedata
import csv
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
import pytesseract
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

# =========================================================
# CONFIGURACOES DE ENTRADA - EDITE SOMENTE ESTA SECAO
# =========================================================
DATA_EVENTO = "01/04/2026"            # DD/MM/AAAA - data base para localizar o evento no HIS/SER
IED         = "NZA 12M7"                      # Nome do IED que sera escrito no arquivo de saida e nos textos
COMENTARIO_BREVE = (
    "Ocorrencia com curto-circuito MONOFÁSICO e atuacao da funcao 51G1T, "
    "seguida de abertura adequada do disjuntor."
)

# =========================================================
# Parametros manuais para calculo do tempo esperado da funcao 51.
# Valores em corrente primaria (A). O script aplica fase para 51P* e neutro para 51G*/51N*.
# =========================================================
PICKUP_FASE_51P = 402
CURVA_FASE_51P = "C2"
DIAL_FASE_51P = "0.20"

PICKUP_NEUTRO_51G = 30
CURVA_NEUTRO_51G = "C2"
DIAL_NEUTRO_51G = "0.90"


# Pasta raiz do projeto. Se preferir, altere para outro caminho absoluto do Windows.
PASTA_PROJETO = Path(__file__).resolve().parent
PASTA_ENTRADA = PASTA_PROJETO / "entrada"
PASTA_MODELO_PPT = PASTA_PROJETO / "MODELO PPT"
PASTA_SAIDA_BASE = PASTA_PROJETO / "saida"

# Se deixar None, o script procurara automaticamente na pasta 'entrada'.
ARQUIVO_TXT = None                     # Ex.: "PCA 12M1__10.75.39.69__02-04-2026__16-36-22.txt"
ARQUIVO_TEMPLATE_PPTX = None           # Ex.: "ANA AT XXX-26 - MODELO.pptx"
ARQUIVO_OSCILOGRAFIA = None            # Ex.: "PCA 12M1 - 1 - SEL-751 - 01_04_2026 - 16_02_23_950 - ABC T.png"
ARQUIVO_OSCILOGRAFIA_CEV = None        # Ex.: "PCA 12M1 - 1 - SEL-751 - 01_04_2026 - 16_02_23_950 - ABC T.CEV"

# Se True, grava um JSON com todos os dados calculados na pasta de saida.
GERAR_JSON_RESUMO = True

# =========================================================
# MODELOS E UTILITARIOS
# =========================================================

SEL_CURVES = {
    "C1": lambda m, td: td * (0.14 / (m ** 0.02 - 1)),
    "C2": lambda m, td: td * (13.5 / (m - 1)),
    "C3": lambda m, td: td * (80 / (m ** 2 - 1)),
    "C4": lambda m, td: td * (120 / (m - 1)),
    "C5": lambda m, td: td * (0.05 / (m ** 0.04 - 1)),
    "U1": lambda m, td: td * (0.0226 + 0.0104 / (m ** 0.02 - 1)),
    "U2": lambda m, td: td * (0.180 + 5.95 / (m ** 2 - 1)),
    "U3": lambda m, td: td * (0.0963 + 3.88 / (m ** 2 - 1)),
    "U4": lambda m, td: td * (0.0352 + 5.67 / (m ** 2 - 1)),
    "U5": lambda m, td: td * (0.00262 + 0.00342 / (m ** 0.02 - 1)),
}

EVENT_TYPE_MAP = {
    "ABC T": "TRIFASICO",
    "AB T": "BIFASICO",
    "BC T": "BIFASICO",
    "CA T": "BIFASICO",
    "AG T": "FASE-TERRA",
    "BG T": "FASE-TERRA",
    "CG T": "FASE-TERRA",
    "ABG T": "BIFASICO-TERRA",
    "BCG T": "BIFASICO-TERRA",
    "CAG T": "BIFASICO-TERRA",
}

SKIP_PROTECTION_NAMES = {
    "TRIP", "FAULT", "ULTRIP", "SV05T", "SV06T", "SV07T", "SV08T", "SV10T",
    "OUT101", "OUT102", "OUT103", "PROTECAO_50", "PROTECAO_51",
}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".webp"}
CEV_EXTS = {".cev"}


@dataclass
class HisRow:
    idx: int
    ref: str
    date: str
    time: str
    event: str
    locat: str
    current: float
    freq: float
    targets: str


@dataclass
class SerRow:
    idx: int
    date: str
    time: str
    element: str
    state: str


# =========================================================
# PARSE DO TXT SEL
# =========================================================


def extract_section(text: str, name: str) -> str:
    pattern = rf"=+\s*{re.escape(name)}\s*=+\n{name}\n(.*?)(?=\n=>|\n=+\s*[A-Z0-9 ]+\s*=+\n[A-Z0-9 ]+\n|\Z)"
    m = re.search(pattern, text, flags=re.S)
    if not m:
        raise ValueError(f"Secao '{name}' nao encontrada no TXT.")
    return m.group(1)


SETTING_RE = re.compile(r"\b([A-Z0-9]+)\s*:=\s*([^\s]+)")


def parse_settings(sho_text: str) -> Dict[str, str]:
    return {k: v for k, v in SETTING_RE.findall(sho_text)}


def parse_ser(ser_text: str) -> List[SerRow]:
    def normalize_relay_date(value: str) -> str:
        value = value.strip()
        if re.match(r"^\d{4}/\d{2}/\d{2}$", value):
            return value
        if re.match(r"^\d{2}/\d{2}/\d{4}$", value):
            return datetime.strptime(value, "%d/%m/%Y").strftime("%Y/%m/%d")
        return value

    rows: List[SerRow] = []
    for line in ser_text.splitlines():
        line = line.rstrip()
        if not re.match(r"^\s*\d+\s+(\d{4}/\d{2}/\d{2}|\d{2}/\d{2}/\d{4})\s+\d{2}:\d{2}:\d{2}\.\d{3}\s+", line):
            continue
        parts = line.split()
        rows.append(SerRow(
            idx=int(parts[0]),
            date=normalize_relay_date(parts[1]),
            time=parts[2],
            element=" ".join(parts[3:-1]),
            state=parts[-1],
        ))
    return rows



def parse_his(his_text: str) -> List[HisRow]:
    def normalize_relay_date(value: str) -> str:
        value = value.strip()
        if re.match(r"^\d{4}/\d{2}/\d{2}$", value):
            return value
        if re.match(r"^\d{2}/\d{2}/\d{4}$", value):
            return datetime.strptime(value, "%d/%m/%Y").strftime("%Y/%m/%d")
        return value

    rows: List[HisRow] = []
    for line in his_text.splitlines():
        line = line.rstrip()
        m = re.match(
            r"^\s*(\d+)\s+(\d+)\s+(\d{4}/\d{2}/\d{2}|\d{2}/\d{2}/\d{4})\s+(\d{2}:\d{2}:\d{2}\.\d{3})\s+(.*?)\s+(\S+)\s+(-?\d+(?:\.\d+)?)\s+(-?\d+(?:\.\d+)?)\s+(\S+)\s*$",
            line,
        )
        if not m:
            continue
        idx, ref, date, time, event, locat, current, freq, targets = m.groups()
        rows.append(HisRow(
            idx=int(idx),
            ref=ref,
            date=normalize_relay_date(date),
            time=time,
            event=event.strip(),
            locat=locat,
            current=float(current),
            freq=float(freq),
            targets=targets,
        ))
    if not rows:
        raise ValueError("Nenhuma linha HIS encontrada.")
    return rows


# =========================================================
# OCR DA OSCILOGRAFIA
# =========================================================


def _ocr_variants(img: Image.Image) -> List[str]:
    variants: List[str] = []
    configs = ["--psm 6", "--psm 11", "--psm 12"]

    gray = ImageOps.grayscale(img)
    candidates = [
        gray,
        ImageOps.autocontrast(gray),
        ImageEnhance.Contrast(ImageOps.autocontrast(gray)).enhance(2.5),
        ImageEnhance.Sharpness(ImageOps.autocontrast(gray)).enhance(2.0),
        ImageOps.invert(ImageEnhance.Contrast(ImageOps.autocontrast(gray)).enhance(2.5)),
    ]

    for base in candidates:
        scaled = base.resize((base.width * 2, base.height * 2))
        scaled = scaled.filter(ImageFilter.SHARPEN)
        for cfg in configs:
            try:
                variants.append(pytesseract.image_to_string(scaled, config=cfg))
            except Exception:
                pass
    return variants



def _normalize_ocr_text(text: str) -> str:
    text = text.replace("|", "I")
    text = text.replace("lA", "IA").replace("1A", "IA")
    text = text.replace("lB", "IB").replace("1B", "IB")
    text = text.replace("lC", "IC").replace("1C", "IC")
    text = text.replace("lG", "IG").replace("1G", "IG")
    text = text.replace("IN:", "IG:")
    text = text.replace("IN ", "IG ")
    text = text.replace("—", "-").replace("−", "-")
    return text



def extract_currents_from_osc_image(image_path: Path) -> Dict[str, float]:
    img = Image.open(image_path)
    w, h = img.size
    right_panel = img.crop((int(w * 0.72), int(h * 0.05), w, int(h * 0.98)))
    texts = _ocr_variants(right_panel)
    combined = "\n".join(_normalize_ocr_text(t) for t in texts)

    m = re.search(
        r"Fault\s*Currents\s*:\s*IA\s*[:=]?\s*(-?\d+[\.,]\d+)\s*IB\s*[:=]?\s*(-?\d+[\.,]\d+)\s*IC\s*[:=]?\s*(-?\d+[\.,]\d+)(?:.*?(?:IG|IN)\s*[:=]?\s*(-?\d+[\.,]\d+))?",
        combined,
        flags=re.I | re.S,
    )
    if m:
        ia, ib, ic, ig = m.groups()
        currents = {
            "IA": float(ia.replace(",", ".")),
            "IB": float(ib.replace(",", ".")),
            "IC": float(ic.replace(",", ".")),
        }
        if ig is not None:
            currents["IG"] = float(ig.replace(",", "."))
        else:
            m_ig = re.search(r"(?:IG|IN)\s*[:=]?\s*(-?\d+[\.,]\d+)", combined, flags=re.I)
            currents["IG"] = float(m_ig.group(1).replace(",", ".")) if m_ig else 0.0
        return currents

    currents: Dict[str, float] = {}
    for label in ["IA", "IB", "IC", "IG", "IN"]:
        mm = re.search(rf"\b{label}\b\s*[:=]?\s*(-?\d+[\.,]\d+)", combined, flags=re.I)
        if mm:
            value = float(mm.group(1).replace(",", "."))
            currents["IG" if label == "IN" else label] = value

    if all(k in currents for k in ["IA", "IB", "IC"]) and "IG" in currents:
        return currents

    raise ValueError(
        "Nao foi possivel extrair IA, IB, IC e IG da oscilografia. "
        "Verifique se o print contem o quadro 'Fault Currents'."
    )


def extract_currents_from_cev_file(cev_path: Path) -> Dict[str, float]:
    with cev_path.open("r", encoding="utf-8", errors="ignore", newline="") as fh:
        rows = list(csv.reader(fh))

    wanted = ("IA(A)", "IB(A)", "IC(A)", "IG(A)")

    for idx, row in enumerate(rows):
        header = [cell.strip().upper() for cell in row]
        if not all(name in header for name in wanted):
            continue

        col_idx = {name: header.index(name) for name in wanted}
        for data_row in rows[idx + 1:]:
            if not data_row:
                continue
            if len(data_row) <= max(col_idx.values()):
                continue
            try:
                ia = float(data_row[col_idx["IA(A)"]].strip().replace(",", "."))
                ib = float(data_row[col_idx["IB(A)"]].strip().replace(",", "."))
                ic = float(data_row[col_idx["IC(A)"]].strip().replace(",", "."))
                ig = float(data_row[col_idx["IG(A)"]].strip().replace(",", "."))
                return {"IA": ia, "IB": ib, "IC": ic, "IG": ig}
            except ValueError:
                continue

    raise ValueError(
        "Nao foi possivel localizar colunas IA(A), IB(A), IC(A), IG(A) com dados validos no arquivo CEV."
    )


def extract_cev_event_location_and_currents(cev_path: Path) -> Dict[str, Any]:
    with cev_path.open("r", encoding="utf-8", errors="ignore", newline="") as fh:
        rows = list(csv.reader(fh))

    wanted = ("IA(A)", "IB(A)", "IC(A)", "IG(A)", "EVENT", "LOCATION")

    for idx, row in enumerate(rows):
        header = [cell.strip().upper() for cell in row]
        if not all(name in header for name in wanted):
            continue

        col_idx = {name: header.index(name) for name in wanted}
        max_idx = max(col_idx.values())

        for data_row in rows[idx + 1:]:
            if not data_row or len(data_row) <= max_idx:
                continue
            try:
                ia = float(data_row[col_idx["IA(A)"]].strip().replace(",", "."))
                ib = float(data_row[col_idx["IB(A)"]].strip().replace(",", "."))
                ic = float(data_row[col_idx["IC(A)"]].strip().replace(",", "."))
                ig = float(data_row[col_idx["IG(A)"]].strip().replace(",", "."))
            except ValueError:
                continue

            event = data_row[col_idx["EVENT"]].strip()
            location_raw = data_row[col_idx["LOCATION"]].strip()
            try:
                location = f"{float(location_raw.replace(',', '.')):.2f}"
            except ValueError:
                location = location_raw

            return {
                "IA": ia,
                "IB": ib,
                "IC": ic,
                "IG": ig,
                "EVENT": event,
                "LOCATION": location,
            }

    raise ValueError(
        "Nao foi possivel localizar colunas IA(A), IB(A), IC(A), IG(A), EVENT e LOCATION com dados validos no arquivo CEV."
    )



def format_currents_text(currents: Dict[str, float]) -> str:
    return (
        f"IA = {currents['IA']:.2f} A ; "
        f"IB = {currents['IB']:.2f} A ; "
        f"IC = {currents['IC']:.2f} A ; "
        f"IG = {currents['IG']:.2f} A"
    )


# =========================================================
# CALCULOS
# =========================================================


def parse_dt(date_str: str, time_str: str) -> datetime:
    return datetime.strptime(f"{date_str} {time_str}", "%Y/%m/%d %H:%M:%S.%f")



def format_hms_ms(delta_seconds: float) -> str:
    total_ms = round(max(delta_seconds, 0.0) * 1000)
    hours, rem = divmod(total_ms, 3_600_000)
    minutes, rem = divmod(rem, 60_000)
    seconds, ms = divmod(rem, 1000)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}.{ms:03d}"



def format_ms(delta_seconds: float) -> str:
    return f"{round(max(delta_seconds, 0.0) * 1000)}ms  (referencia < 100ms)"



def sec_to_human(delta_seconds: float) -> str:
    total_ms = round(max(delta_seconds, 0.0) * 1000)
    hours, rem = divmod(total_ms, 3_600_000)
    minutes, rem = divmod(rem, 60_000)
    seconds, ms = divmod(rem, 1000)
    return f"{hours:02d}h:{minutes:02d}m:{seconds:02d}s.{ms:03d}ms"



def normalize_protection_pick_base(trip_element: str) -> str:
    if not trip_element.endswith("T"):
        raise ValueError(f"Elemento de trip invalido: {trip_element}")
    return trip_element[:-1]



def pick_trip_element(ser_rows: List[SerRow], fault_date: str, trip_time: str) -> str:
    candidates = []
    for row in ser_rows:
        if row.date != fault_date or row.time != trip_time:
            continue
        if row.state not in {"Asserted", "ATUADO"}:
            continue
        if row.element in SKIP_PROTECTION_NAMES:
            continue
        if row.element.startswith(("SV", "OUT", "79", "IN")):
            continue
        if row.element.endswith("T") and row.element[:2] in {"50", "51"}:
            candidates.append(row.element)
    if not candidates:
        raise ValueError("Nao foi possivel identificar a protecao atuada no SER.")
    candidates.sort(key=lambda x: (0 if x.startswith("50P") else 1 if x.startswith("51P") else 2, x))
    return candidates[0]



def find_previous_asserted(ser_rows: List[SerRow], fault_date: str, before_dt: datetime, element: str) -> Optional[SerRow]:
    items = []
    for row in ser_rows:
        if row.date == fault_date and row.element == element and row.state == "Asserted":
            dt = parse_dt(row.date, row.time)
            if dt <= before_dt:
                items.append((dt, row))
    return max(items, default=(None, None), key=lambda x: x[0])[1] if items else None



def find_first_after(ser_rows: List[SerRow], fault_date: str, after_dt: datetime, element: str, state: str) -> Optional[SerRow]:
    items = []
    for row in ser_rows:
        if row.date == fault_date and row.element == element and row.state == state:
            dt = parse_dt(row.date, row.time)
            if dt >= after_dt:
                items.append((dt, row))
    return min(items, default=(None, None), key=lambda x: x[0])[1] if items else None



def calculate_expected_time(
    settings: Dict[str, str],
    protection: str,
    his_current_primary: float,
    cev_currents: Optional[Dict[str, float]] = None,
) -> Optional[float]:
    def parse_num(value: Any) -> float:
        return float(str(value).strip().replace(",", "."))

    def pick_fault_current_for_51_from_cev() -> float:
        if protection.startswith("51P"):
            if cev_currents:
                return max(cev_currents["IA"], cev_currents["IB"], cev_currents["IC"])
            return his_current_primary
        if protection.startswith(("51G", "51N")):
            if cev_currents:
                return cev_currents["IG"]
            return his_current_primary
        return his_current_primary

    base = normalize_protection_pick_base(protection)

    if protection.startswith("50P") or protection.startswith("50G") or protection.startswith("50N"):
        return 0.03

    if protection.startswith("51"):
        if protection.startswith("51P"):
            pickup = parse_num(PICKUP_FASE_51P)
            curve = CURVA_FASE_51P
            dial = parse_num(DIAL_FASE_51P)
            fault_current = pick_fault_current_for_51_from_cev()
        elif protection.startswith(("51G", "51N")):
            pickup = parse_num(PICKUP_NEUTRO_51G)
            curve = CURVA_NEUTRO_51G
            dial = parse_num(DIAL_NEUTRO_51G)
            fault_current = pick_fault_current_for_51_from_cev()
        else:
            ctr = float(settings.get("CTR", "1"))
            fault_current_secondary = his_current_primary / ctr
            pickup = settings.get(f"{base}P")
            curve = settings.get(f"{base}C")
            dial = settings.get(f"{base}TD")
            if not pickup or not curve or not dial:
                return None
            pickup = parse_num(pickup)
            dial = parse_num(dial)
            fault_current = fault_current_secondary

        m = fault_current / pickup
        if m <= 1:
            return None
        fn = SEL_CURVES.get(curve.upper())
        return fn(m, dial) if fn else None

    return None



def infer_fault_type(event: str) -> str:
    return EVENT_TYPE_MAP.get(event.strip().upper(), event.strip().upper())



def sanitize_filename(value: str) -> str:
    value = re.sub(r'[\\/:*?"<>|]', "-", value)
    value = re.sub(r"\s+", " ", value).strip()
    return value



def ensure_input_folder(folder: Path) -> None:
    folder.mkdir(parents=True, exist_ok=True)



def resolve_ied_input_folder(base_folder: Path, ied: str, event_date: str) -> Path:
    expected_name = f"{sanitize_filename(ied)} - {datetime.strptime(event_date, '%d/%m/%Y').strftime('%d.%m.%Y')}"
    candidate = base_folder / expected_name
    if candidate.exists() and candidate.is_dir():
        return candidate
    return base_folder


def select_single_file(folder: Path, explicit_name: Optional[str], suffixes: Tuple[str, ...]) -> Path:
    if explicit_name:
        explicit_path = Path(explicit_name)
        path = explicit_path if explicit_path.is_absolute() else folder / explicit_path
        if not path.exists():
            raise FileNotFoundError(f"Arquivo informado nao encontrado: {path}")
        return path

    files = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in suffixes]
    if not files:
        raise FileNotFoundError(
            f"Nenhum arquivo com extensao {suffixes} foi encontrado em '{folder}'."
        )
    if len(files) > 1:
        names = "\n - ".join(p.name for p in files)
        raise ValueError(
            f"Ha mais de um arquivo candidato em '{folder}'. Defina explicitamente a variavel correspondente no inicio do codigo.\n - {names}"
        )
    return files[0]


def select_optional_single_file(folder: Path, explicit_name: Optional[str], suffixes: Tuple[str, ...]) -> Optional[Path]:
    if explicit_name:
        explicit_path = Path(explicit_name)
        path = explicit_path if explicit_path.is_absolute() else folder / explicit_path
        if not path.exists():
            raise FileNotFoundError(f"Arquivo informado nao encontrado: {path}")
        return path

    files = [p for p in folder.iterdir() if p.is_file() and p.suffix.lower() in suffixes]
    if not files:
        return None
    if len(files) > 1:
        names = "\n - ".join(p.name for p in files)
        raise ValueError(
            f"Ha mais de um arquivo candidato em '{folder}'. Defina explicitamente a variavel correspondente no inicio do codigo.\n - {names}"
        )
    return files[0]



def prepare_output_folder(base_folder: Path, ied: str, event_date: str, execution_dt: datetime) -> Path:
    folder_name = (
        f"ANA AT XXX-26 - {sanitize_filename(ied)} - "
        f"{datetime.strptime(event_date, '%d/%m/%Y').strftime('%d.%m.%Y')} - "
        f"{execution_dt.strftime('%H.%M.%S')}"
    )
    out_dir = base_folder / folder_name
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


def keep_only_final_outputs(out_dir: Path, output_pptx: Path, keep_json: bool) -> None:
    keep_names = {output_pptx.name}
    if keep_json:
        keep_names.add("resumo_dados.json")
    for item in out_dir.iterdir():
        if not item.is_file():
            continue
        if item.name in keep_names:
            continue
        try:
            item.unlink()
        except PermissionError:
            print(f"Aviso: nao foi possivel remover arquivo em uso: {item.name}")



def infer_device_tag(ied: str) -> str:
    match = re.search(r"(\d+[A-Z]\d+)$", ied.replace("_", " ").strip())
    return match.group(1) if match else ied.split()[-1]



def infer_substation_name(ied: str) -> str:
    first = ied.replace("_", " ").split()[0].upper()
    return first



def build_detailed_analysis(data: Dict[str, Any], ied: str, comentario_breve: str) -> str:
    def tone_for_time(actual_ms: int, mech_ms: int) -> str:
        if actual_ms <= 50 and mech_ms <= 80:
            return (
                "Os registros indicam atuacao rapida e coerente com a severidade da falta, sem evidencias de demora indevida "
                "entre a sensibilizacao da protecao e a abertura mecanica do disjuntor."
            )
        if actual_ms <= 150 and mech_ms <= 100:
            return (
                "Os tempos observados permanecem compativeis com uma resposta adequada da protecao e da cadeia de abertura, "
                "sem indicios imediatos de degradacao do desempenho do disjuntor."
            )
        return (
            "Os tempos observados recomendam avaliacao adicional da coordenacao e da cadeia de abertura, com especial atencao "
            "aos circuitos de trip, bobina de abertura e condicoes mecanicas do disjuntor."
        )

    actual_ms = round(data["actual_time_seconds"] * 1000)
    mech_ms = round(data["mechanical_time_seconds"] * 1000)
    expected = data["expected_time"]
    expected_phrase = (
        f"Para a corrente registrada no HIS ({data['his_current']}), o tempo esperado da funcao atuada foi estimado em {expected}. "
        if expected != "N/D" else
        "Nao foi possivel estimar com seguranca o tempo esperado a partir dos ajustes disponiveis, motivo pelo qual o comparativo teorico deve ser tratado como nao disponivel. "
    )
    locator_phrase = (
        f"O localizador de falta estava habilitado e apontou distancia estimada de {data['fault_distance']}, informacao util para correlacao com a rede e com os vestigios de campo. "
        if data["fault_locator"] == "SIM" else
        "O localizador de falta nao estava habilitado para esta ocorrencia, nao havendo distancia estimada registrada no relatorio. "
    )
    reclose_phrase = (
        "O religamento automatico encontrava-se habilitado nos ajustes do rele. "
        if data["reclose"].startswith("SIM") else
        "O religamento automatico nao estava habilitado para esta condicao operacional. "
    )

    return (
        f"Na ocorrencia em analise, o IED {ied} registrou um evento do tipo {data['fault_type'].lower()}, com corrente de falta de {data['his_current']} no HIS e atuacao da funcao {data['protection']} as {data['trip_time']}. "
        f"No SER, observa-se a sensibilizacao inicial em {data['pickup_time']}, resultando em tempo real de atuacao do rele de {data['actual_time']} ({actual_ms} ms). "
        f"A abertura mecanica do disjuntor foi confirmada em {data['open_time']}, correspondendo a {data['mechanical_time']} ({mech_ms} ms). "
        f"{expected_phrase}"
        f"{tone_for_time(actual_ms, mech_ms)} "
        f"{reclose_phrase}"
        f"{locator_phrase}"
        f"Comentario complementar do analista: {comentario_breve.strip()}"
    )


# =========================================================
# POWERPOINT
# =========================================================


def replace_text(shape, old: str, new: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if old not in shape.text:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)



def replace_regex_in_shape(shape, pattern: str, replacement: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    original = shape.text
    updated = re.sub(pattern, replacement, original, flags=re.I)
    if updated != original:
        shape.text = updated



def set_table_value(table_shape, label: str, new_value: str) -> bool:
    def normalize_label(text: str) -> str:
        clean = unicodedata.normalize("NFKD", text)
        clean = "".join(ch for ch in clean if not unicodedata.combining(ch))
        clean = re.sub(r"\s+", " ", clean).strip().upper()
        return clean

    target = normalize_label(label)
    table = table_shape.table
    for r in range(len(table.rows)):
        current = normalize_label(table.cell(r, 0).text)
        if current == target:
            table.cell(r, 1).text = new_value
            return True
    return False



def find_main_table(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape
    raise ValueError("Tabela de analise nao encontrada no slide.")


def find_slide_with_table(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                return slide
    raise ValueError("Tabela de analise nao encontrada no template PPTX.")



def find_large_text_box(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            candidates.append((shape.width * shape.height, shape))
    return max(candidates, default=(None, None), key=lambda x: x[0])[1]



def apply_text_frame_font(text_frame, font_name: str, font_size_pt: float) -> None:
    for paragraph in text_frame.paragraphs:
        if not paragraph.runs:
            run = paragraph.add_run()
            run.text = paragraph.text or ""
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = RGBColor(0, 0, 0)


def set_shape_text_with_font(shape, text: str, font_name: str, font_size_pt: float) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = RGBColor(0, 0, 0)


def style_table_cells_aptos(table_shape, font_name: str = "Aptos", font_size_pt: float = 11) -> None:
    table = table_shape.table
    for r in range(len(table.rows)):
        for c in range(len(table.columns)):
            cell = table.cell(r, c)
            text_frame = cell.text_frame
            for paragraph in text_frame.paragraphs:
                if not paragraph.runs:
                    run = paragraph.add_run()
                    run.text = paragraph.text or ""
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(font_size_pt)
                    # Coluna da esquerda (rotulos) em branco; resultados em preto.
                    run.font.color.rgb = RGBColor(255, 255, 255) if c == 0 else RGBColor(0, 0, 0)


def _norm_text(value: str) -> str:
    clean = unicodedata.normalize("NFKD", value)
    clean = "".join(ch for ch in clean if not unicodedata.combining(ch))
    clean = re.sub(r"\s+", " ", clean).strip().lower()
    return clean


def find_comment_and_image_slots(slide):
    heading = None
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if "analise de desempenho da protecao" in _norm_text(shape.text or ""):
            heading = shape
            break

    if heading is None:
        return None, None

    comment_candidates = []
    image_candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or shape is heading:
            continue
        text_norm = _norm_text(shape.text or "")
        if text_norm.startswith("desarme do dj") or text_norm.startswith("registro de eventos"):
            continue
        if shape.top > heading.top:
            comment_candidates.append(shape)
        elif shape.top < heading.top:
            image_candidates.append(shape)

    comment_shape = max(comment_candidates, key=lambda s: s.width * s.height) if comment_candidates else None
    image_shape = max(image_candidates, key=lambda s: s.width * s.height) if image_candidates else None
    return comment_shape, image_shape


def replace_picture_in_slot(slide, image_path: str, slot_shape, table_shape=None) -> bool:
    if slot_shape is None:
        return False
    left, top, width, height = slot_shape.left, slot_shape.top, slot_shape.width, slot_shape.height
    if table_shape is not None:
        margin = 100000
        # Alinha imagem com a janela de dados (mesmo topo/altura da tabela).
        top = table_shape.top
        height = table_shape.height
        max_right = table_shape.left - margin
        if max_right > left:
            width = max_right - left
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    old = slot_shape._element
    old.getparent().remove(old)
    return True



def update_presentation(template_path: Path, output_path: Path, data: Dict[str, Any], osc_image: Optional[Path], ied: str) -> Optional[str]:
    prs = Presentation(str(template_path))
    device_tag = infer_device_tag(ied)

    if len(prs.slides) < 2:
        raise ValueError("O template deve possuir pelo menos 2 slides para preencher o slide 2.")

    # Regra do template: preencher somente o slide 2, mantendo os demais sem alteracoes.
    target_slide = prs.slides[1]
    for shape in target_slide.shapes:
        replace_text(shape, "01/04/2026", data["occurrence_date"])
        replace_regex_in_shape(shape, r"ANA AT XXX-26\s*-\s*.*", f"ANA AT XXX-26 - {ied}")

    table_shape = find_main_table(target_slide)
    set_table_value(table_shape, "TIPO DE DEFEITO:", data["fault_type"])
    set_table_value(table_shape, "CORRENTE DE CARGA:", data["current_load"])
    set_table_value(table_shape, "PROTECAO ATUADA:", data["protection"])
    set_table_value(table_shape, "HORARIO DO DISPARO (RELE)", data["trip_time"])
    set_table_value(table_shape, "RELIGAMENTO AUTOMATICO:", data["reclose"])
    set_table_value(table_shape, "TEMPO DE ATUACAO REAL (RELE):", data["actual_time"])
    set_table_value(table_shape, "TEMPO DE ATUACAO ESPERADO:", data["expected_time"])
    set_table_value(table_shape, "TEMPO DE RESPOSTA MECANICA:", data["mechanical_time"])
    set_table_value(table_shape, "LOCALIZADOR DE FALTA:", data["fault_locator"])
    set_table_value(table_shape, "DISTANCIA REAL DO DEFEITO:", data["fault_distance"])

    # compatibilidade com template acentuado
    for label, value in [
        ("PROTEÇÃO ATUADA:", data["protection"]),
        ("HORÁRIO DO DISPARO (RELÉ)", data["trip_time"]),
        ("RELIGAMENTO AUTOMÁTICO:", data["reclose"]),
        ("TEMPO DE ATUAÇÃO REAL (RELÉ):", data["actual_time"]),
        ("TEMPO DE ATUAÇÃO ESPERADO:", data["expected_time"]),
        ("TEMPO DE RESPOSTA MECÂNICA:", data["mechanical_time"]),
        ("LOCALIZADOR DE FALTA:", data["fault_locator"]),
        ("DISTÂNCIA REAL DO DEFEITO:", data["fault_distance"]),
    ]:
        set_table_value(table_shape, label, value)
    style_table_cells_aptos(table_shape, "Aptos", 11)

    comment_shape, image_slot = find_comment_and_image_slots(target_slide)
    if comment_shape is not None:
        set_shape_text_with_font(comment_shape, data["detailed_commentary"], "Aptos", 11)

    for shape in target_slide.shapes:
        if getattr(shape, "has_text_frame", False):
            replace_regex_in_shape(shape, r"DESARME DO DJ\s+[A-Z0-9]+", f"DESARME DO DJ {device_tag}")
            replace_regex_in_shape(shape, r"SE\s+[A-Z0-9_]+\s*\([A-Z0-9_]+\)", f"SE {infer_substation_name(ied)} ({infer_substation_name(ied)})")

    if osc_image and osc_image.exists():
        replace_picture_in_slot(target_slide, str(osc_image), image_slot, table_shape)

    prs.save(str(output_path))
    return None


# =========================================================
# PIPELINE
# =========================================================


def build_data(
    txt_path: Path,
    osc_image: Optional[Path],
    osc_cev: Optional[Path],
    occurrence_date: str,
    ied: str,
    comentario_breve: str,
) -> Dict[str, Any]:
    text = txt_path.read_text(encoding="utf-8", errors="ignore")
    sho_text = extract_section(text, "SHO")
    ser_text = extract_section(text, "SER")
    his_text = extract_section(text, "HIS")

    settings = parse_settings(sho_text)
    ser_rows = parse_ser(ser_text)
    his_rows = parse_his(his_text)

    occ = datetime.strptime(occurrence_date, "%d/%m/%Y").strftime("%Y/%m/%d")
    matches = [r for r in his_rows if r.date == occ]
    if not matches:
        raise ValueError(f"Nenhum evento HIS encontrado para {occurrence_date}.")
    his_row = matches[0]

    fault_date = his_row.date
    trip_time = his_row.time
    protection = pick_trip_element(ser_rows, fault_date, trip_time)
    pickup_element = normalize_protection_pick_base(protection) + "P"

    trip_dt = parse_dt(fault_date, trip_time)
    pickup_row = find_previous_asserted(ser_rows, fault_date, trip_dt, pickup_element)
    if not pickup_row:
        raise ValueError(f"Pickup inicial '{pickup_element}' nao encontrado no SER.")
    pickup_dt = parse_dt(pickup_row.date, pickup_row.time)

    open_row = find_first_after(ser_rows, fault_date, trip_dt, "DISJUNTOR", "ABERTO")
    if not open_row:
        raise ValueError("Evento 'DISJUNTOR ABERTO' nao encontrado apos o trip.")
    open_dt = parse_dt(open_row.date, open_row.time)

    actual_time = (trip_dt - pickup_dt).total_seconds()
    mechanical_time = (open_dt - trip_dt).total_seconds()

    reclose = "SIM" if settings.get("E79", "OFF") != "OFF" else "NA (SEM RELIGAMENTO)"
    efloc_enabled = settings.get("EFLOC", "N") == "Y"

    currents_dict = None
    cev_currents_dict: Optional[Dict[str, float]] = None
    cev_event: Optional[str] = None
    cev_location: Optional[str] = None
    current_warning: Optional[str] = None
    current_source = "N/D"
    if osc_cev:
        try:
            cev_data = extract_cev_event_location_and_currents(osc_cev)
            currents_dict = {
                "IA": cev_data["IA"],
                "IB": cev_data["IB"],
                "IC": cev_data["IC"],
                "IG": cev_data["IG"],
            }
            cev_currents_dict = dict(currents_dict)
            cev_event = cev_data["EVENT"]
            cev_location = cev_data["LOCATION"]
            current_source = "CEV"
        except Exception as exc:
            current_warning = f"Aviso: falha ao extrair correntes do CEV ({osc_cev.name}). Detalhe: {exc}"

    if currents_dict is None and osc_image:
        try:
            currents_dict = extract_currents_from_osc_image(osc_image)
            current_source = "OCR_IMAGEM"
        except Exception as exc:
            ocr_msg = (
                "Aviso: nao foi possivel extrair correntes da oscilografia via OCR. "
                f"Sera usado 'N/D' para corrente de carga. Detalhe: {exc}"
            )
            current_warning = f"{current_warning} | {ocr_msg}" if current_warning else ocr_msg
    current_load = format_currents_text(currents_dict) if currents_dict else "N/D"

    fault_event_source = cev_event if cev_event else his_row.event
    fault_distance = cev_location if cev_location is not None else (his_row.locat if efloc_enabled else "NA")
    fault_locator = "SIM" if (cev_location is not None or efloc_enabled) else "NA"

    expected_time = calculate_expected_time(
        settings=settings,
        protection=protection,
        his_current_primary=his_row.current,
        cev_currents=cev_currents_dict,
    )

    data: Dict[str, Any] = {
        "ied": ied,
        "occurrence_date": datetime.strptime(fault_date, "%Y/%m/%d").strftime("%d/%m/%Y"),
        "fault_type": infer_fault_type(fault_event_source),
        "current_load": current_load,
        "protection": protection,
        "trip_time": trip_time,
        "pickup_time": pickup_row.time,
        "open_time": open_row.time,
        "reclose": reclose,
        "actual_time": format_hms_ms(actual_time),
        "expected_time": format_hms_ms(expected_time) if expected_time is not None else "N/D",
        "mechanical_time": format_ms(mechanical_time),
        "fault_locator": fault_locator,
        "fault_distance": fault_distance,
        "his_current": f"{his_row.current:.1f} A",
        "his_event": his_row.event,
        "ctr": settings.get("CTR", "1"),
        "actual_time_seconds": actual_time,
        "mechanical_time_seconds": mechanical_time,
        "expected_time_seconds": expected_time,
        "comentario_breve": comentario_breve.strip(),
        "corrente_fonte": current_source,
        "ocr_warning": current_warning,
    }
    if currents_dict:
        data.update({
            "osc_IA": f"{currents_dict['IA']:.2f}",
            "osc_IB": f"{currents_dict['IB']:.2f}",
            "osc_IC": f"{currents_dict['IC']:.2f}",
            "osc_IG": f"{currents_dict['IG']:.2f}",
        })

    data["detailed_commentary"] = build_detailed_analysis(data, ied, comentario_breve)
    return data



def main() -> None:
    ensure_input_folder(PASTA_ENTRADA)
    PASTA_MODELO_PPT.mkdir(parents=True, exist_ok=True)
    PASTA_SAIDA_BASE.mkdir(parents=True, exist_ok=True)
    execution_dt = datetime.now()

    input_folder = resolve_ied_input_folder(PASTA_ENTRADA, IED, DATA_EVENTO)

    txt_path = select_single_file(input_folder, ARQUIVO_TXT, (".txt",))
    template_path = select_single_file(PASTA_MODELO_PPT, ARQUIVO_TEMPLATE_PPTX, (".pptx",))
    osc_path = select_optional_single_file(input_folder, ARQUIVO_OSCILOGRAFIA, tuple(IMAGE_EXTS))
    cev_path = select_optional_single_file(input_folder, ARQUIVO_OSCILOGRAFIA_CEV, tuple(CEV_EXTS))

    out_dir = prepare_output_folder(PASTA_SAIDA_BASE, IED, DATA_EVENTO, execution_dt)
    output_pptx = out_dir / f"ANA AT XXX-26 - {sanitize_filename(IED)} - {datetime.strptime(DATA_EVENTO, '%d/%m/%Y').strftime('%d.%m.%Y')}.pptx"

    data = build_data(
        txt_path=txt_path,
        osc_image=osc_path,
        osc_cev=cev_path,
        occurrence_date=DATA_EVENTO,
        ied=IED,
        comentario_breve=COMENTARIO_BREVE,
    )
    template_warning = update_presentation(template_path, output_pptx, data, osc_path, IED)

    if GERAR_JSON_RESUMO:
        json_path = out_dir / "resumo_dados.json"
        json_ready = {k: v for k, v in data.items() if not k.endswith("_seconds")}
        json_path.write_text(json.dumps(json_ready, indent=2, ensure_ascii=False), encoding="utf-8")

    keep_only_final_outputs(out_dir, output_pptx, GERAR_JSON_RESUMO)

    if data.get("ocr_warning"):
        print(f"\n{data['ocr_warning']}")
    if template_warning:
        print(f"\n{template_warning}")

    print("Arquivo gerado com sucesso:")
    print(output_pptx)
    print("\nPasta de saida:")
    print(out_dir)


if __name__ == "__main__":
    main()
