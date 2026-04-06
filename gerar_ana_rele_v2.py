#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
import pytesseract
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# =========================
# MODELOS E UTILITÁRIOS
# =========================

SEL_CURVES = {
    "C1": lambda m, td: td * (0.14 / (m ** 0.02 - 1)),
    "C2": lambda m, td: td * (13.5 / (m - 1)),
    "C3": lambda m, td: td * (80 / (m ** 2 - 1)),
    "C4": lambda m, td: td * (120 / (m - 1)),
    "C5": lambda m, td: td * (0.05 / (m ** 0.04 - 1)),
    "U1": lambda m, td: td * (0.0226 + 0.0104 / (m ** 0.02 - 1)),
    "U2": lambda m, td: td * (0.180 + 5.95 / (m ** 2 - 1)),
    "U3": lambda m, td: td * (0.0963 + 3.88 / (m ** 2 - 1)),
    "U4": lambda m, td: td * (0.0352 + 5.67 / (m ** 2 - 1)),
    "U5": lambda m, td: td * (0.00262 + 0.00342 / (m ** 0.02 - 1)),
}

EVENT_TYPE_MAP = {
    "ABC T": "TRIFÁSICO",
    "AB T": "BIFÁSICO",
    "BC T": "BIFÁSICO",
    "CA T": "BIFÁSICO",
    "AG T": "FASE-TERRA",
    "BG T": "FASE-TERRA",
    "CG T": "FASE-TERRA",
    "ABG T": "BIFÁSICO-TERRA",
    "BCG T": "BIFÁSICO-TERRA",
    "CAG T": "BIFÁSICO-TERRA",
}

SKIP_PROTECTION_NAMES = {
    "TRIP", "FAULT", "ULTRIP", "SV05T", "SV06T", "SV07T", "SV08T", "SV10T",
    "OUT101", "OUT102", "OUT103", "PROTECAO_50", "PROTECAO_51",
}


@dataclass
class HisRow:
    idx: int
    ref: str
    date: str
    time: str
    event: str
    locat: str
    current: float
    freq: float
    targets: str


@dataclass
class SerRow:
    idx: int
    date: str
    time: str
    element: str
    state: str


# =========================
# PARSE DO TXT SEL
# =========================


def extract_section(text: str, name: str) -> str:
    pattern = rf"=+\s*{re.escape(name)}\s*=+\n{name}\n(.*?)(?=\n=>|\n=+\s*[A-Z0-9 ]+\s*=+\n[A-Z0-9 ]+\n|\Z)"
    m = re.search(pattern, text, flags=re.S)
    if not m:
        raise ValueError(f"Seção '{name}' não encontrada no TXT.")
    return m.group(1)


SETTING_RE = re.compile(r"\b([A-Z0-9]+)\s*:=\s*([^\s]+)")


def parse_settings(sho_text: str) -> Dict[str, str]:
    return {k: v for k, v in SETTING_RE.findall(sho_text)}


def parse_ser(ser_text: str) -> List[SerRow]:
    rows: List[SerRow] = []
    for line in ser_text.splitlines():
        line = line.rstrip()
        if not re.match(r"^\s*\d+\s+\d{4}/\d{2}/\d{2}\s+\d{2}:\d{2}:\d{2}\.\d{3}\s+", line):
            continue
        parts = line.split()
        rows.append(SerRow(
            idx=int(parts[0]),
            date=parts[1],
            time=parts[2],
            element=" ".join(parts[3:-1]),
            state=parts[-1],
        ))
    return rows


def parse_his(his_text: str) -> List[HisRow]:
    rows: List[HisRow] = []
    for line in his_text.splitlines():
        line = line.rstrip()
        if not re.match(r"^\s*\d+\s+\d+\s+\d{4}/\d{2}/\d{2}\s+\d{2}:\d{2}:\d{2}\.\d{3}\s+", line):
            continue
        parts = line.split()
        rows.append(HisRow(
            idx=int(parts[0]),
            ref=parts[1],
            date=parts[2],
            time=parts[3],
            event=" ".join(parts[4:-4]),
            locat=parts[-4],
            current=float(parts[-3]),
            freq=float(parts[-2]),
            targets=parts[-1],
        ))
    if not rows:
        raise ValueError("Nenhuma linha HIS encontrada.")
    return rows


# =========================
# OCR DA OSCILOGRAFIA
# =========================


def _ocr_variants(img: Image.Image) -> List[str]:
    variants = []
    configs = ["--psm 6", "--psm 11", "--psm 12"]

    candidates = []
    gray = ImageOps.grayscale(img)
    candidates.append(gray)
    candidates.append(ImageOps.autocontrast(gray))
    candidates.append(ImageEnhance.Contrast(ImageOps.autocontrast(gray)).enhance(2.5))
    candidates.append(ImageEnhance.Sharpness(ImageOps.autocontrast(gray)).enhance(2.0))
    candidates.append(ImageOps.invert(ImageEnhance.Contrast(ImageOps.autocontrast(gray)).enhance(2.5)))

    for base in candidates:
        scaled = base.resize((base.width * 2, base.height * 2))
        scaled = scaled.filter(ImageFilter.SHARPEN)
        for cfg in configs:
            try:
                variants.append(pytesseract.image_to_string(scaled, config=cfg))
            except Exception:
                pass
    return variants


def _normalize_ocr_text(text: str) -> str:
    text = text.replace("|", "I")
    text = text.replace("lA", "IA").replace("1A", "IA")
    text = text.replace("lB", "IB").replace("1B", "IB")
    text = text.replace("lC", "IC").replace("1C", "IC")
    text = text.replace("lG", "IG").replace("1G", "IG")
    text = text.replace("IN:", "IG:")
    text = text.replace("IN ", "IG ")
    text = text.replace("—", "-").replace("−", "-")
    return text


def extract_currents_from_osc_image(image_path: Path) -> Dict[str, float]:
    img = Image.open(image_path)
    w, h = img.size

    # área da lateral direita onde costuma aparecer "Fault Currents"
    right_panel = img.crop((int(w * 0.72), int(h * 0.05), w, int(h * 0.98)))

    texts = _ocr_variants(right_panel)
    combined = "\n".join(_normalize_ocr_text(t) for t in texts)

    # tenta primeiro localizar a linha "Fault Currents"
    m = re.search(
        r"Fault\s*Currents\s*:\s*IA\s*[:=]?\s*(-?\d+[\.,]\d+)\s*IB\s*[:=]?\s*(-?\d+[\.,]\d+)\s*IC\s*[:=]?\s*(-?\d+[\.,]\d+)(?:.*?(?:IG|IN)\s*[:=]?\s*(-?\d+[\.,]\d+))?",
        combined,
        flags=re.I | re.S,
    )
    if m:
        ia, ib, ic, ig = m.groups()
        currents = {
            "IA": float(ia.replace(",", ".")),
            "IB": float(ib.replace(",", ".")),
            "IC": float(ic.replace(",", ".")),
        }
        if ig is not None:
            currents["IG"] = float(ig.replace(",", "."))
        else:
            # fallback: procura IG/IN fora da linha principal
            m_ig = re.search(r"(?:IG|IN)\s*[:=]?\s*(-?\d+[\.,]\d+)", combined, flags=re.I)
            currents["IG"] = float(m_ig.group(1).replace(",", ".")) if m_ig else 0.0
        return currents

    # fallback: extrai individualmente
    currents: Dict[str, float] = {}
    for label in ["IA", "IB", "IC", "IG", "IN"]:
        mm = re.search(rf"\b{label}\b\s*[:=]?\s*(-?\d+[\.,]\d+)", combined, flags=re.I)
        if mm:
            value = float(mm.group(1).replace(",", "."))
            currents["IG" if label == "IN" else label] = value

    if all(k in currents for k in ["IA", "IB", "IC"]) and "IG" in currents:
        return currents

    raise ValueError(
        "Não foi possível extrair IA, IB, IC e IG da oscilografia. "
        "Verifique se o print contém o quadro 'Fault Currents'."
    )


def format_currents_text(currents: Dict[str, float]) -> str:
    return (
        f"IA = {currents['IA']:.2f} A ; "
        f"IB = {currents['IB']:.2f} A ; "
        f"IC = {currents['IC']:.2f} A ; "
        f"IG = {currents['IG']:.2f} A"
    )


# =========================
# CÁLCULOS
# =========================


def parse_dt(date_str: str, time_str: str) -> datetime:
    return datetime.strptime(f"{date_str} {time_str}", "%Y/%m/%d %H:%M:%S.%f")


def format_hms_ms(delta_seconds: float) -> str:
    total_ms = round(max(delta_seconds, 0.0) * 1000)
    hours, rem = divmod(total_ms, 3_600_000)
    minutes, rem = divmod(rem, 60_000)
    seconds, ms = divmod(rem, 1000)
    return f"{hours:02d}:{minutes:02d}:{seconds:02d}.{ms:03d}"


def format_ms(delta_seconds: float) -> str:
    return f"{round(max(delta_seconds, 0.0) * 1000)}ms  (referência < 100ms)"


def sec_to_human(delta_seconds: float) -> str:
    total_ms = round(max(delta_seconds, 0.0) * 1000)
    hours, rem = divmod(total_ms, 3_600_000)
    minutes, rem = divmod(rem, 60_000)
    seconds, ms = divmod(rem, 1000)
    return f"{hours:02d}h:{minutes:02d}m:{seconds:02d}s.{ms:03d}ms"


def normalize_protection_pick_base(trip_element: str) -> str:
    if not trip_element.endswith("T"):
        raise ValueError(f"Elemento de trip inválido: {trip_element}")
    return trip_element[:-1]


def pick_trip_element(ser_rows: List[SerRow], fault_date: str, trip_time: str) -> str:
    candidates = []
    for row in ser_rows:
        if row.date != fault_date or row.time != trip_time:
            continue
        if row.state not in {"Asserted", "ATUADO"}:
            continue
        if row.element in SKIP_PROTECTION_NAMES:
            continue
        if row.element.startswith(("SV", "OUT", "79", "IN")):
            continue
        if row.element.endswith("T") and row.element[:2] in {"50", "51"}:
            candidates.append(row.element)
    if not candidates:
        raise ValueError("Não foi possível identificar a proteção atuada no SER.")
    candidates.sort(key=lambda x: (0 if x.startswith("50P") else 1 if x.startswith("51P") else 2, x))
    return candidates[0]


def find_previous_asserted(ser_rows: List[SerRow], fault_date: str, before_dt: datetime, element: str) -> Optional[SerRow]:
    items = []
    for row in ser_rows:
        if row.date == fault_date and row.element == element and row.state == "Asserted":
            dt = parse_dt(row.date, row.time)
            if dt <= before_dt:
                items.append((dt, row))
    return max(items, default=(None, None), key=lambda x: x[0])[1] if items else None


def find_first_after(ser_rows: List[SerRow], fault_date: str, after_dt: datetime, element: str, state: str) -> Optional[SerRow]:
    items = []
    for row in ser_rows:
        if row.date == fault_date and row.element == element and row.state == state:
            dt = parse_dt(row.date, row.time)
            if dt >= after_dt:
                items.append((dt, row))
    return min(items, default=(None, None), key=lambda x: x[0])[1] if items else None


def calculate_expected_time(settings: Dict[str, str], protection: str, his_current_primary: float) -> Optional[float]:
    base = normalize_protection_pick_base(protection)
    ctr = float(settings.get("CTR", "1"))
    fault_current_secondary = his_current_primary / ctr

    if protection.startswith("50"):
        delay = settings.get(f"{base}D")
        return float(delay) if delay and delay != "OFF" else None

    if protection.startswith("51"):
        pickup = settings.get(f"{base}P")
        curve = settings.get(f"{base}C")
        dial = settings.get(f"{base}TD")
        if not pickup or not curve or not dial:
            return None
        m = fault_current_secondary / float(pickup)
        if m <= 1:
            return None
        fn = SEL_CURVES.get(curve.upper())
        return fn(m, float(dial)) if fn else None

    return None


def infer_fault_type(event: str) -> str:
    return EVENT_TYPE_MAP.get(event.strip().upper(), event.strip().upper())


# =========================
# POWERPOINT
# =========================


def replace_text(shape, old: str, new: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    if old not in shape.text:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def set_table_value(table_shape, label: str, new_value: str) -> bool:
    table = table_shape.table
    for r in range(len(table.rows)):
        if table.cell(r, 0).text.strip() == label:
            table.cell(r, 1).text = new_value
            return True
    return False


def find_main_table(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_table", False):
            return shape
    raise ValueError("Tabela de análise não encontrada no slide 2.")


def find_large_text_box(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and len(shape.text or "") > 80:
            candidates.append((shape.width * shape.height, shape))
    return max(candidates, default=(None, None), key=lambda x: x[0])[1]


def replace_picture(slide, image_path: str) -> bool:
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pics:
        return False
    pic = max(pics, key=lambda s: s.width * s.height)
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    old = pic._element
    old.getparent().remove(old)
    return True


def build_summary_text(data: Dict[str, str]) -> str:
    return (
        f"O disjuntor PCA_12M1 registrou um {data['fault_type'].lower()} com atuação da função "
        f"{data['protection']} às {data['trip_time']} após pickup em {data['pickup_time']}, "
        f"resultando em {data['actual_time_human']} de tempo de atuação do relé. "
        f"A abertura mecânica ocorreu às {data['open_time']}, correspondendo a {data['mechanical_time_human']} "
        f"de tempo de resposta mecânica do disjuntor."
    )


# =========================
# PIPELINE
# =========================


def build_data(
    txt_path: Path,
    corrente_carga: Optional[str] = None,
    occurrence_date: Optional[str] = None,
    osc_image: Optional[Path] = None,
) -> Dict[str, str]:
    text = txt_path.read_text(encoding="utf-8", errors="ignore")
    sho_text = extract_section(text, "SHO")
    ser_text = extract_section(text, "SER")
    his_text = extract_section(text, "HIS")

    settings = parse_settings(sho_text)
    ser_rows = parse_ser(ser_text)
    his_rows = parse_his(his_text)

    his_row = his_rows[0]
    if occurrence_date:
        occ = datetime.strptime(occurrence_date, "%d/%m/%Y").strftime("%Y/%m/%d")
        matches = [r for r in his_rows if r.date == occ]
        if not matches:
            raise ValueError(f"Nenhum evento HIS encontrado para {occurrence_date}.")
        his_row = matches[0]

    fault_date = his_row.date
    trip_time = his_row.time
    protection = pick_trip_element(ser_rows, fault_date, trip_time)
    pickup_element = normalize_protection_pick_base(protection) + "P"

    trip_dt = parse_dt(fault_date, trip_time)
    pickup_row = find_previous_asserted(ser_rows, fault_date, trip_dt, pickup_element)
    if not pickup_row:
        raise ValueError(f"Pickup inicial '{pickup_element}' não encontrado no SER.")
    pickup_dt = parse_dt(pickup_row.date, pickup_row.time)

    open_row = find_first_after(ser_rows, fault_date, trip_dt, "DISJUNTOR", "ABERTO")
    if not open_row:
        raise ValueError("Evento 'DISJUNTOR ABERTO' não encontrado após o trip.")
    open_dt = parse_dt(open_row.date, open_row.time)

    actual_time = (trip_dt - pickup_dt).total_seconds()
    mechanical_time = (open_dt - trip_dt).total_seconds()
    expected_time = calculate_expected_time(settings, protection, his_row.current)

    reclose = "SIM" if settings.get("E79", "OFF") != "OFF" else "NA (SEM RELIGAMENTO)"
    efloc_enabled = settings.get("EFLOC", "N") == "Y"
    fault_distance = his_row.locat if efloc_enabled else "NA"

    currents_dict = None
    if osc_image:
        currents_dict = extract_currents_from_osc_image(osc_image)
        corrente_carga = format_currents_text(currents_dict)
    elif not corrente_carga:
        corrente_carga = "N/D"

    data = {
        "occurrence_date": datetime.strptime(fault_date, "%Y/%m/%d").strftime("%d/%m/%Y"),
        "fault_type": infer_fault_type(his_row.event),
        "current_load": corrente_carga,
        "protection": protection,
        "trip_time": trip_time,
        "pickup_time": pickup_row.time,
        "open_time": open_row.time,
        "reclose": reclose,
        "actual_time": format_hms_ms(actual_time),
        "actual_time_human": sec_to_human(actual_time),
        "expected_time": format_hms_ms(expected_time) if expected_time is not None else "N/D",
        "mechanical_time": format_ms(mechanical_time),
        "mechanical_time_human": f"{round(mechanical_time * 1000)} ms",
        "fault_locator": "SIM" if efloc_enabled else "NA",
        "fault_distance": fault_distance,
        "his_current": f"{his_row.current:.1f} A",
        "his_event": his_row.event,
        "ctr": settings.get("CTR", "1"),
        "expected_time_sec": "" if expected_time is None else f"{expected_time:.6f}",
    }

    if currents_dict:
        data["osc_IA"] = f"{currents_dict['IA']:.2f}"
        data["osc_IB"] = f"{currents_dict['IB']:.2f}"
        data["osc_IC"] = f"{currents_dict['IC']:.2f}"
        data["osc_IG"] = f"{currents_dict['IG']:.2f}"

    return data


def update_presentation(template_path: Path, output_path: Path, data: Dict[str, str], osc_image: Optional[Path] = None) -> None:
    prs = Presentation(str(template_path))

    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        replace_text(shape, "01/04/2026", data["occurrence_date"])

    slide2 = prs.slides[1]
    table_shape = find_main_table(slide2)
    set_table_value(table_shape, "TIPO DE DEFEITO:", data["fault_type"])
    set_table_value(table_shape, "CORRENTE DE CARGA:", data["current_load"])
    set_table_value(table_shape, "PROTEÇÃO ATUADA:", data["protection"])
    set_table_value(table_shape, "HORÁRIO DO DISPARO (RELÉ)", data["trip_time"])
    set_table_value(table_shape, "RELIGAMENTO AUTOMÁTICO:", data["reclose"])
    set_table_value(table_shape, "TEMPO DE ATUAÇÃO REAL (RELÉ):", data["actual_time"])
    set_table_value(table_shape, "TEMPO DE ATUAÇÃO ESPERADO:", data["expected_time"])
    set_table_value(table_shape, "TEMPO DE RESPOSTA MECÂNICA:", data["mechanical_time"])
    set_table_value(table_shape, "LOCALIZADOR DE FALTA:", data["fault_locator"])
    set_table_value(table_shape, "DISTÂNCIA REAL DO DEFEITO:", data["fault_distance"])

    body = find_large_text_box(slide2)
    if body is not None:
        body.text = build_summary_text(data)

    if osc_image and osc_image.exists():
        replace_picture(slide2, str(osc_image))

    prs.save(str(output_path))


# =========================
# CLI
# =========================


def main() -> None:
    parser = argparse.ArgumentParser(description="Gera um ANA em PowerPoint a partir do template, TXT SEL e oscilografia.")
    parser.add_argument("--template", required=True, help="Arquivo .pptx modelo")
    parser.add_argument("--txt", required=True, help="Arquivo .txt coletado do relé")
    parser.add_argument("--output", required=True, help="Arquivo .pptx de saída")
    parser.add_argument("--osc", help="Imagem da oscilografia/print principal")
    parser.add_argument("--corrente-carga", help="Sobrescreve o campo CORRENTE DE CARGA")
    parser.add_argument("--occurrence-date", help="Data da ocorrência em DD/MM/AAAA; se omitido, usa a 1ª linha do HIS")
    parser.add_argument("--dump-json", action="store_true", help="Mostra no terminal os dados calculados")
    args = parser.parse_args()

    osc_path = Path(args.osc) if args.osc else None
    data = build_data(
        Path(args.txt),
        corrente_carga=args.corrente_carga,
        occurrence_date=args.occurrence_date,
        osc_image=osc_path,
    )
    update_presentation(Path(args.template), Path(args.output), data, osc_path)

    if args.dump_json:
        print(json.dumps(data, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
